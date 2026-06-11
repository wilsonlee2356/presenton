import asyncio
import os
from pathlib import Path
import re
import tempfile
import urllib
import zipfile
from functools import lru_cache
import aiohttp
import xml.etree.ElementTree as ET
from typing import Dict, List, Optional, Sequence, Set, Tuple
from pydantic import BaseModel
from pptx import Presentation
from fontTools.ttLib import TTFont

DEFAULT_GOOGLE_FONT_WEIGHTS = (400, 700)
PPT_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_TEXT_STYLE_TAGS = ("a:rPr", "a:defRPr", "a:endParaRPr")
_FONT_TAGS = ("a:latin", "a:ea", "a:cs")
_THEME_FONT_REFERENCES = {"+mn-lt", "+mj-lt", "+mn-ea", "+mj-ea", "+mn-cs", "+mj-cs"}
_SFNT_FORMATS = {
    "\x00\x01\x00\x00": "TrueType",
    "true": "TrueType",
    "typ1": "PostScript Type 1",
    "OTTO": "OpenType CFF",
    "ttcf": "TrueType Collection",
    "wOFF": "WOFF",
    "wOF2": "WOFF2",
}


class FontDetail(BaseModel):
    """Detailed information about a font file."""

    file: str
    size_bytes: int
    error: Optional[str] = None
    eot_extraction_error: Optional[str] = None
    family_name: Optional[str] = None
    subfamily_name: Optional[str] = None
    unique_id: Optional[str] = None
    full_name: Optional[str] = None
    version: Optional[str] = None
    postscript_name: Optional[str] = None
    trademark: Optional[str] = None
    manufacturer: Optional[str] = None
    designer: Optional[str] = None
    description: Optional[str] = None
    vendor_url: Optional[str] = None
    designer_url: Optional[str] = None
    license: Optional[str] = None
    license_url: Optional[str] = None
    weight_class: Optional[int] = None
    width_class: Optional[int] = None
    cap_height: Optional[int] = None
    x_height: Optional[int] = None
    ascent: Optional[int] = None
    descent: Optional[int] = None
    units_per_em: Optional[int] = None
    created: Optional[int] = None
    modified: Optional[int] = None
    ascender: Optional[int] = None
    descender: Optional[int] = None
    line_gap: Optional[int] = None
    num_glyphs: Optional[int] = None
    format: Optional[str] = None


_STYLE_TOKENS = {
    "italic",
    "italics",
    "ital",
    "oblique",
    "roman",
    "gras",
    "bolditalic",
    "bolditalics",
    "thin",
    "hairline",
    "extralight",
    "ultralight",
    "light",
    "demilight",
    "semilight",
    "book",
    "regular",
    "normal",
    "medium",
    "semibold",
    "demibold",
    "bold",
    "extrabold",
    "ultrabold",
    "black",
    "extrablack",
    "ultrablack",
    "heavy",
    "narrow",
    "condensed",
    "semicondensed",
    "extracondensed",
    "ultracondensed",
    "expanded",
    "semiexpanded",
    "extraexpanded",
    "ultraexpanded",
}
_STYLE_MODIFIERS = {"semi", "demi", "extra", "ultra"}


def _clean_font_metadata_string(value: str) -> str:
    return "".join(
        char
        for char in value
        if char == "\t" or char == "\n" or char == "\r" or ord(char) >= 32
    ).strip()


def _normalize_font_format(value: object) -> Optional[str]:
    if not value:
        return None
    if isinstance(value, bytes):
        raw = value.decode("latin1", errors="ignore")
    else:
        raw = str(value)
    return _SFNT_FORMATS.get(raw, _clean_font_metadata_string(raw) or None)


def normalize_font_family_name(raw_name: str) -> str:
    """Normalize raw font family labels by trimming weight/style tokens."""
    if not raw_name:
        return raw_name
    name = raw_name.replace("_", " ").replace("-", " ")
    name = re.sub(r"(?<=[a-z0-9])([A-Z])", r" \1", name)
    name = re.sub(r"([A-Z]+)([A-Z][a-z])", r"\1 \2", name)
    name = re.sub(r"\s+", " ", name).strip()
    lower_name = name.lower()
    for style in sorted(_STYLE_TOKENS, key=len, reverse=True):
        if lower_name.endswith(" " + style):
            name = name[: -(len(style) + 1)]
            lower_name = lower_name[: -(len(style) + 1)]
            break
    tokens_original = name.split(" ")
    tokens_filtered: List[str] = []
    for index, tok in enumerate(tokens_original):
        lower_tok = tok.lower()
        if index == 0:
            tokens_filtered.append(tok)
            continue
        if lower_tok in _STYLE_TOKENS or lower_tok in _STYLE_MODIFIERS:
            continue
        tokens_filtered.append(tok)
    if not tokens_filtered:
        tokens_filtered = tokens_original
    normalized = " ".join(tokens_filtered).strip()
    normalized = re.sub(r"\s+", " ", normalized)
    return normalized


def build_google_fonts_stylesheet_url(
    family_name: str,
    weights: Optional[Sequence[int]] = DEFAULT_GOOGLE_FONT_WEIGHTS,
    variants: Optional[Sequence[str]] = None,
) -> str:
    encoded_family = urllib.parse.quote_plus(family_name)
    requested_variants = set(variants or [])
    requested_weights = set(weights or [])
    if requested_variants:
        requested_weights = {400}
        if "bold" in requested_variants or "bold_italic" in requested_variants:
            requested_weights.add(700)
    if requested_weights:
        normalized_weights = sorted(
            {int(weight) for weight in requested_weights if int(weight) > 0}
        )
        if "italic" in requested_variants or "bold_italic" in requested_variants:
            italic_weights = set()
            if "italic" in requested_variants:
                italic_weights.add(400)
            if "bold_italic" in requested_variants:
                italic_weights.add(700)
            weights_param = ";".join(
                [*(f"0,{weight}" for weight in normalized_weights)]
                + [*(f"1,{weight}" for weight in sorted(italic_weights))]
            )
            return (
                "https://fonts.googleapis.com/css2"
                f"?family={encoded_family}:ital,wght@{weights_param}&display=swap"
            )
        weight_selector = ";".join(str(weight) for weight in normalized_weights)
        return (
            "https://fonts.googleapis.com/css2"
            f"?family={encoded_family}:wght@{weight_selector}&display=swap"
        )
    return f"https://fonts.googleapis.com/css2?family={encoded_family}&display=swap"


def _resolve_theme_typeface(
    typeface: Optional[str], theme_fonts: Optional[Dict[str, str]] = None
) -> Optional[str]:
    if not typeface:
        return None
    cleaned = typeface.strip()
    if not cleaned:
        return None
    if cleaned.startswith("+mj"):
        resolved = (theme_fonts or {}).get("major", "")
        return resolved.strip() or None
    if cleaned.startswith("+mn"):
        resolved = (theme_fonts or {}).get("minor", "")
        return resolved.strip() or None
    return cleaned


def _extract_typefaces_from_text_style_node(
    text_style_node: ET.Element, theme_fonts: Optional[Dict[str, str]] = None
) -> List[str]:
    fonts: List[str] = []
    seen = set()
    font_tags = _FONT_TAGS
    latin_elem = text_style_node.find("a:latin", PPT_NS)
    latin_typeface = (
        _resolve_theme_typeface(latin_elem.get("typeface"), theme_fonts)
        if latin_elem is not None
        else None
    )
    if latin_typeface and latin_typeface not in _THEME_FONT_REFERENCES:
        # Latin text can carry ea/cs fallback faces; counting those adds false positives.
        font_tags = ("a:latin",)
    for font_tag in font_tags:
        font_elem = text_style_node.find(font_tag, PPT_NS)
        if font_elem is None:
            continue
        resolved = _resolve_theme_typeface(font_elem.get("typeface"), theme_fonts)
        if not resolved or resolved in _THEME_FONT_REFERENCES or resolved in seen:
            continue
        seen.add(resolved)
        fonts.append(resolved)
    return fonts


def _extract_fonts_from_xml_root(
    root: ET.Element, theme_fonts: Optional[Dict[str, str]] = None
) -> Set[str]:
    fonts: Set[str] = set()
    for style_tag in _TEXT_STYLE_TAGS:
        for style_elem in root.findall(f".//{style_tag}", PPT_NS):
            fonts.update(_extract_typefaces_from_text_style_node(style_elem, theme_fonts))
    return fonts


def extract_fonts_from_oxml(xml_content: str) -> List[str]:
    """Extract font names referenced by text style nodes inside PPTX XML."""
    try:
        root = ET.fromstring(xml_content)
        return sorted(_extract_fonts_from_xml_root(root))
    except Exception as exc:
        print(f"Error extracting fonts from OXML: {exc}")
        return []


# Helper: Fetch TTF/OTF file URLs for a Google Fonts family via Webfonts Developer API
async def get_google_font_file_urls(family_name: str, api_key: str) -> List[str]:
    encoded_family = urllib.parse.quote_plus(family_name)
    api_url = f"https://www.googleapis.com/webfonts/v1/webfonts?family={encoded_family}&key={api_key}"
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(
                api_url, timeout=aiohttp.ClientTimeout(total=20)
            ) as resp:
                if resp.status != 200:
                    return []
                data = await resp.json()
                items = data.get("items", []) or []
                if not items:
                    return []
                urls: List[str] = []
                # Take first matching family
                files = (items[0] or {}).get("files", {}) or {}
                for _variant, url in files.items():
                    if not url:
                        continue
                    # Prefer directly loadable TTF/OTF files; upgrade to https.
                    fixed_url = url.replace("http://", "https://")
                    lower = fixed_url.lower()
                    if lower.endswith(".ttf") or lower.endswith(".otf"):
                        urls.append(fixed_url)
                return urls
    except Exception:
        return []


async def check_google_font_availability(
    font_name: str, variants: Optional[Sequence[str]] = None
) -> bool:
    """Return True when Google Fonts serves the requested family/variants."""
    try:
        url = build_google_fonts_stylesheet_url(font_name, variants=variants)
        async with aiohttp.ClientSession() as session:
            async with session.get(
                url, timeout=aiohttp.ClientTimeout(total=10)
            ) as response:
                if response.status != 200:
                    return False
                css = await response.text()
                if "@font-face" not in css:
                    return False
                return "fonts.gstatic.com/l/font?kit=" not in css
    except Exception as exc:
        print(f"Error checking Google Font availability for {font_name}: {exc}")
        return False


def extract_raw_fonts_and_embedded_details(
    pptx_path: str, temp_dir: str
) -> Tuple[Set[str], List[FontDetail], List[str]]:
    raw_fonts = extract_used_fonts_from_pptx(pptx_path)

    emb_font_details: List[FontDetail] = []
    emb_font_paths: List[str] = []

    if not raw_fonts:
        return raw_fonts, emb_font_details, emb_font_paths

    with zipfile.ZipFile(pptx_path, "r") as zip_ref:
        emb_fonts_paths_rel = [
            path
            for path in zip_ref.namelist()
            if path.startswith("ppt/fonts/") and path.endswith(".fntdata")
        ]

        for rel_path in emb_fonts_paths_rel:
            try:
                zip_ref.extract(rel_path, temp_dir)
                font_path = os.path.join(temp_dir, rel_path)

                detail = get_font_details(font_path)
                emb_font_details.append(detail)
                emb_font_paths.append(font_path)

            except zipfile.BadZipFile:
                # CORRUPTED EMBEDDED FONT — SKIP
                print(f"Skipping corrupted embedded font: {rel_path}")
                continue
            except Exception as exc:
                # Font parsing failed — also skip
                print(f"Failed to parse embedded font {rel_path}: {exc}")
                continue

    return raw_fonts, emb_font_details, emb_font_paths


def normalize_font_variants(variants: Optional[Sequence[str]]) -> List[str]:
    order = ("regular", "bold", "italic", "bold_italic")
    variant_set = set(variants or [])
    if not variant_set:
        variant_set = {"regular"}
    return [variant for variant in order if variant in variant_set]


def _merge_font_variants(
    target: Dict[str, Set[str]], source: Dict[str, Set[str]]
) -> None:
    for font_name, variants in source.items():
        target.setdefault(font_name, set()).update(variants)


def _is_truthy_ooxml_flag(value: Optional[str]) -> bool:
    return str(value or "").lower() in {"1", "true", "on"}


def _font_style_variant(
    font_name: str,
    r_pr: Optional[ET.Element],
    default_rprs: Sequence[ET.Element] = (),
) -> str:
    bold = False
    italic = False
    for style_node in [r_pr, *default_rprs]:
        if style_node is None:
            continue
        if style_node.get("b") is not None:
            bold = _is_truthy_ooxml_flag(style_node.get("b"))
            break
    for style_node in [r_pr, *default_rprs]:
        if style_node is None:
            continue
        if style_node.get("i") is not None:
            italic = _is_truthy_ooxml_flag(style_node.get("i"))
            break

    inferred_weight = _extract_weight_from_name(font_name)
    if inferred_weight == "bold":
        bold = True
    compact_name = _normalize_compact(font_name)
    if "italic" in compact_name or "oblique" in compact_name:
        italic = True

    if bold and italic:
        return "bold_italic"
    if bold:
        return "bold"
    if italic:
        return "italic"
    return "regular"


def extract_used_font_variants_from_pptx(pptx_path: str) -> Dict[str, Set[str]]:
    """Return font names and regular/bold/italic variants used by slide content."""

    def _local_name(tag: str) -> str:
        if "}" in tag:
            return tag.rsplit("}", 1)[-1]
        return tag

    def _read_zip_xml(zip_ref: zipfile.ZipFile, path: str) -> Optional[ET.Element]:
        try:
            return ET.fromstring(zip_ref.read(path))
        except Exception:
            return None

    def _get_relationships(
        zip_ref: zipfile.ZipFile, path: str
    ) -> Dict[str, Dict[str, str]]:
        dir_name = os.path.dirname(path)
        filename = os.path.basename(path)
        rels_path = os.path.join(dir_name, "_rels", f"{filename}.rels").replace("\\", "/")
        rels_xml = _read_zip_xml(zip_ref, rels_path)
        if rels_xml is None:
            return {}

        rels: Dict[str, Dict[str, str]] = {}
        for rel in rels_xml.findall(f"{{{REL_NS}}}Relationship"):
            rel_id = rel.get("Id")
            rel_type = rel.get("Type")
            target = rel.get("Target")
            if not rel_id or not rel_type or not target:
                continue
            if target.startswith("/"):
                resolved = target[1:]
            else:
                resolved = os.path.normpath(os.path.join(dir_name, target)).replace(
                    "\\", "/"
                )
            rels[rel_id] = {"path": resolved, "type": rel_type}
        return rels

    def _load_theme_fonts(zip_ref: zipfile.ZipFile) -> Dict[str, str]:
        presentation_xml = _read_zip_xml(zip_ref, "ppt/presentation.xml")
        if presentation_xml is None:
            return {}
        pres_rels = _get_relationships(zip_ref, "ppt/presentation.xml")
        theme_path = next(
            (
                rel["path"]
                for rel in pres_rels.values()
                if "theme" in rel.get("type", "")
            ),
            "ppt/theme/theme1.xml",
        )
        theme_xml = _read_zip_xml(zip_ref, theme_path)
        if theme_xml is None:
            return {}

        font_scheme = theme_xml.find(".//a:fontScheme", PPT_NS)
        if font_scheme is None:
            return {}

        theme_fonts: Dict[str, str] = {}
        major = font_scheme.find("a:majorFont/a:latin", PPT_NS)
        minor = font_scheme.find("a:minorFont/a:latin", PPT_NS)
        if major is not None and major.get("typeface"):
            theme_fonts["major"] = major.get("typeface", "").strip()
        if minor is not None and minor.get("typeface"):
            theme_fonts["minor"] = minor.get("typeface", "").strip()
        return theme_fonts

    def _get_slide_paths(zip_ref: zipfile.ZipFile) -> List[str]:
        presentation_xml = _read_zip_xml(zip_ref, "ppt/presentation.xml")
        pres_rels = _get_relationships(zip_ref, "ppt/presentation.xml")
        if presentation_xml is None:
            slide_paths = [
                name
                for name in zip_ref.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            ]
            slide_paths.sort(
                key=lambda name: int(
                    os.path.basename(name).replace("slide", "").replace(".xml", "")
                )
            )
            return slide_paths

        slide_id_list = presentation_xml.find("p:sldIdLst", PPT_NS)
        if slide_id_list is None:
            return []

        slide_paths: List[str] = []
        rel_attr = f"{{{PPT_NS['r']}}}id"
        for slide_id in slide_id_list.findall("p:sldId", PPT_NS):
            rel_id = slide_id.get(rel_attr)
            rel_info = pres_rels.get(rel_id or "")
            if not rel_info or "slide" not in rel_info.get("type", ""):
                continue
            slide_paths.append(rel_info["path"])
        return slide_paths

    def _is_placeholder(shape: ET.Element) -> bool:
        nv_sp_pr = shape.find("p:nvSpPr", PPT_NS)
        if nv_sp_pr is None:
            return False
        nv_pr = nv_sp_pr.find("p:nvPr", PPT_NS)
        if nv_pr is None:
            return False
        return nv_pr.find("p:ph", PPT_NS) is not None

    def _is_hidden(shape: ET.Element) -> bool:
        nv_container = shape.find("p:nvSpPr", PPT_NS)
        if nv_container is None:
            nv_container = shape.find("p:nvPicPr", PPT_NS)
        if nv_container is None:
            nv_container = shape.find("p:nvGraphicFramePr", PPT_NS)
        if nv_container is None:
            return False
        c_nv_pr = nv_container.find("p:cNvPr", PPT_NS)
        if c_nv_pr is None:
            return False
        return c_nv_pr.get("hidden") in {"1", "true"}

    def _get_placeholder_key(shape: ET.Element) -> Optional[Tuple[str, Optional[str]]]:
        nv_sp_pr = shape.find("p:nvSpPr", PPT_NS)
        if nv_sp_pr is None:
            return None
        nv_pr = nv_sp_pr.find("p:nvPr", PPT_NS)
        if nv_pr is None:
            return None
        ph = nv_pr.find("p:ph", PPT_NS)
        if ph is None:
            return None
        return (ph.get("type") or "body", ph.get("idx"))

    def _placeholder_style_key(ph_type: str) -> str:
        if ph_type in {"title", "ctrTitle"}:
            return "title"
        if ph_type == "body":
            return "body"
        return "other"

    def _build_placeholder_text_style_map(
        layout_xml: Optional[ET.Element], master_xml: Optional[ET.Element]
    ) -> Dict[Tuple[str, Optional[str]], Dict[int, List[ET.Element]]]:
        style_map: Dict[Tuple[str, Optional[str]], Dict[int, List[ET.Element]]] = {}

        tx_styles = master_xml.find("p:txStyles", PPT_NS) if master_xml is not None else None
        txstyle_defaults: Dict[str, Dict[int, ET.Element]] = {}
        if tx_styles is not None:
            for name, key in (
                ("p:titleStyle", "title"),
                ("p:bodyStyle", "body"),
                ("p:otherStyle", "other"),
            ):
                style_elem = tx_styles.find(name, PPT_NS)
                if style_elem is None:
                    continue
                per_level: Dict[int, ET.Element] = {}
                for level in range(1, 10):
                    lvl_pr = style_elem.find(f"a:lvl{level}pPr", PPT_NS)
                    if lvl_pr is None:
                        continue
                    def_rpr = lvl_pr.find("a:defRPr", PPT_NS)
                    if def_rpr is not None:
                        per_level[level - 1] = def_rpr
                if per_level:
                    txstyle_defaults[key] = per_level

        for xml_root in (master_xml, layout_xml):
            if xml_root is None:
                continue
            sp_tree = xml_root.find(".//p:spTree", PPT_NS)
            if sp_tree is None:
                continue
            for child in sp_tree:
                if _local_name(child.tag) != "sp":
                    continue
                placeholder_key = _get_placeholder_key(child)
                if not placeholder_key:
                    continue
                base_defaults = txstyle_defaults.get(
                    _placeholder_style_key(placeholder_key[0]), {}
                )
                tx_body = child.find("p:txBody", PPT_NS)
                lst_style = (
                    tx_body.find("a:lstStyle", PPT_NS) if tx_body is not None else None
                )
                per_level: Dict[int, List[ET.Element]] = {}
                for level in range(1, 10):
                    defaults: List[ET.Element] = []
                    if lst_style is not None:
                        lvl_pr = lst_style.find(f"a:lvl{level}pPr", PPT_NS)
                        if lvl_pr is not None:
                            def_rpr = lvl_pr.find("a:defRPr", PPT_NS)
                            if def_rpr is not None:
                                defaults.append(def_rpr)
                    if level - 1 in base_defaults:
                        defaults.append(base_defaults[level - 1])
                    if defaults:
                        per_level[level - 1] = defaults
                if per_level:
                    style_map[placeholder_key] = per_level
        return style_map

    def _paragraph_level(p_pr: Optional[ET.Element]) -> int:
        if p_pr is None:
            return 0
        level_value = p_pr.get("lvl")
        if level_value is None:
            return 0
        try:
            return int(level_value)
        except ValueError:
            return 0

    def _build_local_text_style_map(tx_body: ET.Element) -> Dict[int, List[ET.Element]]:
        lst_style = tx_body.find("a:lstStyle", PPT_NS)
        if lst_style is None:
            return {}
        local_map: Dict[int, List[ET.Element]] = {}
        for level in range(1, 10):
            lvl_pr = lst_style.find(f"a:lvl{level}pPr", PPT_NS)
            if lvl_pr is None:
                continue
            def_rpr = lvl_pr.find("a:defRPr", PPT_NS)
            if def_rpr is not None:
                local_map[level - 1] = [def_rpr]
        return local_map

    def _get_default_rprs(
        p_pr: Optional[ET.Element],
        local_text_styles: Dict[int, List[ET.Element]],
        placeholder_text_styles: Optional[
            Dict[Tuple[str, Optional[str]], Dict[int, List[ET.Element]]]
        ],
        placeholder_key: Optional[Tuple[str, Optional[str]]],
    ) -> List[ET.Element]:
        level = _paragraph_level(p_pr)
        defaults: List[ET.Element] = list(local_text_styles.get(level, []))
        if not placeholder_text_styles or not placeholder_key:
            return defaults
        style_map = placeholder_text_styles.get(placeholder_key)
        if style_map is None and placeholder_key[0]:
            style_map = placeholder_text_styles.get((placeholder_key[0], None))
        if style_map:
            defaults.extend(style_map.get(level, []))
        return defaults

    def _extract_effective_run_font_variants(
        r_pr: Optional[ET.Element],
        default_rprs: Sequence[ET.Element],
        theme_fonts: Dict[str, str],
    ) -> Dict[str, Set[str]]:
        variant_fonts: Dict[str, Set[str]] = {}
        if r_pr is not None:
            direct_fonts = _extract_typefaces_from_text_style_node(r_pr, theme_fonts)
            if direct_fonts:
                for font_name in direct_fonts:
                    variant_fonts.setdefault(font_name, set()).add(
                        _font_style_variant(font_name, r_pr, default_rprs)
                    )
                return variant_fonts
        for default_rpr in default_rprs:
            inherited_fonts = _extract_typefaces_from_text_style_node(
                default_rpr, theme_fonts
            )
            if inherited_fonts:
                for font_name in inherited_fonts:
                    variant_fonts.setdefault(font_name, set()).add(
                        _font_style_variant(font_name, r_pr, [default_rpr])
                    )
                return variant_fonts
        return variant_fonts

    def _collect_fonts_from_text_body(
        tx_body: ET.Element,
        placeholder_text_styles: Optional[
            Dict[Tuple[str, Optional[str]], Dict[int, List[ET.Element]]]
        ],
        placeholder_key: Optional[Tuple[str, Optional[str]]],
        theme_fonts: Dict[str, str],
    ) -> Dict[str, Set[str]]:
        font_variants: Dict[str, Set[str]] = {}
        local_text_styles = _build_local_text_style_map(tx_body)
        run_tags = {f"{{{PPT_NS['a']}}}r", f"{{{PPT_NS['a']}}}fld"}

        for paragraph in tx_body.findall("a:p", PPT_NS):
            p_pr = paragraph.find("a:pPr", PPT_NS)
            default_rprs = _get_default_rprs(
                p_pr, local_text_styles, placeholder_text_styles, placeholder_key
            )
            for child in paragraph:
                if child.tag not in run_tags:
                    continue
                text_node = child.find("a:t", PPT_NS)
                if text_node is None or not text_node.text:
                    continue
                _merge_font_variants(
                    font_variants,
                    _extract_effective_run_font_variants(
                        child.find("a:rPr", PPT_NS),
                        default_rprs,
                        theme_fonts,
                    ),
                )
        return font_variants

    def _iter_shape_nodes(parent: ET.Element):
        for child in parent:
            tag_name = _local_name(child.tag)
            if tag_name == "grpSp":
                yield from _iter_shape_nodes(child)
                continue
            if tag_name in {"sp", "graphicFrame"}:
                yield child

    def _collect_fonts_from_shape_tree(
        sp_tree: ET.Element,
        theme_fonts: Dict[str, str],
        skip_placeholders: bool = False,
        placeholder_text_styles: Optional[
            Dict[Tuple[str, Optional[str]], Dict[int, List[ET.Element]]]
        ] = None,
    ) -> Dict[str, Set[str]]:
        font_variants: Dict[str, Set[str]] = {}
        for shape in _iter_shape_nodes(sp_tree):
            if _is_hidden(shape):
                continue
            if _local_name(shape.tag) == "sp":
                if skip_placeholders and _is_placeholder(shape):
                    continue
                tx_body = shape.find("p:txBody", PPT_NS)
                if tx_body is None:
                    continue
                _merge_font_variants(
                    font_variants,
                    _collect_fonts_from_text_body(
                        tx_body,
                        placeholder_text_styles,
                        _get_placeholder_key(shape),
                        theme_fonts,
                    ),
                )
                continue
            for tx_body in shape.findall(".//a:txBody", PPT_NS):
                _merge_font_variants(
                    font_variants,
                    _collect_fonts_from_text_body(
                        tx_body,
                        placeholder_text_styles=None,
                        placeholder_key=None,
                        theme_fonts=theme_fonts,
                    ),
                )
        return font_variants

    raw_font_variants: Dict[str, Set[str]] = {}
    try:
        with zipfile.ZipFile(pptx_path, "r") as zip_ref:
            theme_fonts = _load_theme_fonts(zip_ref)

            for slide_path in _get_slide_paths(zip_ref):
                slide_xml = _read_zip_xml(zip_ref, slide_path)
                if slide_xml is None:
                    continue
                slide_rels = _get_relationships(zip_ref, slide_path)
                layout_path = next(
                    (
                        rel["path"]
                        for rel in slide_rels.values()
                        if "slideLayout" in rel.get("type", "")
                    ),
                    None,
                )

                layout_xml = _read_zip_xml(zip_ref, layout_path) if layout_path else None
                layout_rels = (
                    _get_relationships(zip_ref, layout_path) if layout_path else {}
                )
                master_path = next(
                    (
                        rel["path"]
                        for rel in layout_rels.values()
                        if "slideMaster" in rel.get("type", "")
                    ),
                    None,
                )
                master_xml = _read_zip_xml(zip_ref, master_path) if master_path else None
                placeholder_text_styles = _build_placeholder_text_style_map(
                    layout_xml, master_xml
                )

                if master_xml is not None:
                    master_sp_tree = master_xml.find(".//p:spTree", PPT_NS)
                    if master_sp_tree is not None:
                        _merge_font_variants(
                            raw_font_variants,
                            _collect_fonts_from_shape_tree(
                                master_sp_tree,
                                theme_fonts=theme_fonts,
                                skip_placeholders=True,
                            ),
                        )

                if layout_xml is not None:
                    layout_sp_tree = layout_xml.find(".//p:spTree", PPT_NS)
                    if layout_sp_tree is not None:
                        _merge_font_variants(
                            raw_font_variants,
                            _collect_fonts_from_shape_tree(
                                layout_sp_tree,
                                theme_fonts=theme_fonts,
                                skip_placeholders=True,
                            ),
                        )

                slide_sp_tree = slide_xml.find(".//p:spTree", PPT_NS)
                if slide_sp_tree is not None:
                    _merge_font_variants(
                        raw_font_variants,
                        _collect_fonts_from_shape_tree(
                            slide_sp_tree,
                            theme_fonts=theme_fonts,
                            placeholder_text_styles=placeholder_text_styles,
                        ),
                    )

            for name in zip_ref.namelist():
                if not name.startswith("ppt/charts/") or not name.endswith(".xml"):
                    continue
                chart_xml = _read_zip_xml(zip_ref, name)
                if chart_xml is None:
                    continue
                for font_name in _extract_fonts_from_xml_root(chart_xml, theme_fonts):
                    raw_font_variants.setdefault(font_name, set()).add("regular")
    except Exception:
        print("Failed to read PPTX XML parts, returning empty fonts list")
        return {}

    return raw_font_variants


def extract_used_fonts_from_pptx(pptx_path: str) -> Set[str]:
    """Return all font names referenced in a PPTX (slides, masters, layouts, theme)."""
    return set(extract_used_font_variants_from_pptx(pptx_path).keys())


async def get_available_and_unavailable_fonts_for_pptx(
    pptx_path: str, temp_dir: str
) -> Tuple[List[Tuple[str, Optional[str]]], List[Tuple[str, Optional[str]]]]:
    """
    Return lists of available/unavailable fonts for a PPTX file.

    Args:
        pptx_path: Path to the PPTX file to inspect.
        temp_dir: Temporary directory for extracted assets.

    Returns:
        Tuple of (available_fonts, unavailable_fonts) where each entry is a list
        of (font_name, url or None).
    """
    raw_fonts, emb_font_details, _ = await asyncio.to_thread(
        extract_raw_fonts_and_embedded_details,
        pptx_path,
        temp_dir,
    )
    font_variants_by_name = await asyncio.to_thread(
        extract_used_font_variants_from_pptx, pptx_path
    )

    if not raw_fonts:
        return [], []

    found_fonts_with_url: Dict[str, str] = {}
    for font_name in raw_fonts:
        match_index = get_index_of_matching_font_detail_or_none(
            font_name, emb_font_details
        )
        if match_index is None:
            continue
        found_fonts_with_url[font_name] = (
            "https://example.com/just-a-placeholder-url.ttf"
        )

    matched_fonts = set(found_fonts_with_url.keys())
    fonts_to_check = sorted(raw_fonts - matched_fonts)

    normalized_variants: Dict[str, Set[str]] = {}
    for font_name, variants in font_variants_by_name.items():
        normalized_name = normalize_font_family_name(font_name)
        if normalized_name:
            normalized_variants.setdefault(normalized_name, set()).update(variants)

    fonts_to_check = [normalize_font_family_name(font) for font in fonts_to_check]
    fonts_to_check = list(set(fonts_to_check))

    availability_results: List[bool] = []
    if fonts_to_check:
        availability_results = await asyncio.gather(
            *[
                check_google_font_availability(
                    font,
                    variants=normalize_font_variants(normalized_variants.get(font)),
                )
                for font in fonts_to_check
            ]
        )

    available_fonts: List[Tuple[str, Optional[str]]] = []
    unavailable_fonts: List[Tuple[str, Optional[str]]] = []

    for font_name, font_url in found_fonts_with_url.items():
        available_fonts.append((font_name, font_url))

    for font, is_available in zip(fonts_to_check, availability_results):
        if is_available:
            google_fonts_url = build_google_fonts_stylesheet_url(
                font,
                variants=normalize_font_variants(normalized_variants.get(font)),
            )
            available_fonts.append((font, google_fonts_url))
        else:
            unavailable_fonts.append((font, None))

    return available_fonts, unavailable_fonts


def create_font_alias_config(
    raw_fonts: List[str],
    extra_includes: Optional[List[str]] = None,
    temp_dir: Optional[str] = None,
    explicit_aliases: Optional[Dict[str, str]] = None,
    protected_font_names: Optional[Sequence[str]] = None,
) -> str:
    """Create a fontconfig alias file mapping variant families to normalized names."""
    mappings: Dict[str, str] = {}
    explicit_aliases = {
        src: dst
        for src, dst in (explicit_aliases or {}).items()
        if src and dst and src != dst
    }
    protected_names = {name for name in (protected_font_names or []) if name}
    explicit_names = set(explicit_aliases.keys()).union(explicit_aliases.values())
    skip_normalization = protected_names.union(explicit_names)
    for font_name in raw_fonts:
        if font_name in skip_normalization:
            continue
        normalized = normalize_font_family_name(font_name)
        if normalized and normalized != font_name:
            mappings[font_name] = normalized
    fd, fonts_conf_path = tempfile.mkstemp(
        prefix="fonts_alias_",
        suffix=".conf",
        dir=temp_dir,
    )
    os.close(fd)
    with open(fonts_conf_path, "w", encoding="utf-8") as cfg:
        cfg.write(
            """<?xml version='1.0'?>
<!DOCTYPE fontconfig SYSTEM "urn:fontconfig:fonts.dtd">
<fontconfig>
  <include>/etc/fonts/fonts.conf</include>
"""
        )
        if extra_includes:
            for include_path in extra_includes:
                if not include_path:
                    continue
                cfg.write(f"  <include>{include_path}</include>\n")
        for src, dst in explicit_aliases.items():
            cfg.write(
                f"""
  <match target="pattern">
    <test name="family" compare="eq">
      <string>{src}</string>
    </test>
    <edit name="family" mode="assign" binding="strong">
      <string>{dst}</string>
    </edit>
  </match>
"""
            )
        for src, dst in mappings.items():
            cfg.write(
                f"""
  <match target="pattern">
    <test name="family" compare="eq">
      <string>{src}</string>
    </test>
    <edit name="family" mode="assign" binding="strong">
      <string>{dst}</string>
    </edit>
  </match>
"""
            )
        cfg.write("\n</fontconfig>\n")
    return fonts_conf_path


def _replace_fonts_in_xml_root(
    root: ET.Element,
    font_mapping: Dict[str, str],
    font_variant_mapping: Optional[Dict[str, Dict[str, str]]] = None,
) -> bool:
    def _first_typeface(style_elem: Optional[ET.Element]) -> Optional[str]:
        if style_elem is None:
            return None
        for font_tag in _FONT_TAGS:
            font_elem = style_elem.find(font_tag, PPT_NS)
            if font_elem is not None and font_elem.get("typeface"):
                return font_elem.get("typeface")
        return None

    changed = False
    for style_tag in _TEXT_STYLE_TAGS:
        for style_elem in root.findall(f".//{style_tag}", PPT_NS):
            for font_tag in _FONT_TAGS:
                font_elem = style_elem.find(font_tag, PPT_NS)
                if font_elem is None:
                    continue
                typeface = font_elem.get("typeface")
                if not typeface:
                    continue
                replacement = None
                variant_mapping = (font_variant_mapping or {}).get(typeface)
                if variant_mapping:
                    variant = _font_style_variant(typeface, style_elem, [])
                    replacement = variant_mapping.get(variant)
                if replacement is None:
                    replacement = font_mapping.get(typeface)
                if replacement and replacement != typeface:
                    font_elem.set("typeface", replacement)
                    changed = True

    run_tags = {f"{{{PPT_NS['a']}}}r", f"{{{PPT_NS['a']}}}fld"}
    for paragraph in root.findall(".//a:p", PPT_NS):
        p_pr = paragraph.find("a:pPr", PPT_NS)
        paragraph_default = (
            p_pr.find("a:defRPr", PPT_NS) if p_pr is not None else None
        )
        inherited_typeface = _first_typeface(paragraph_default)
        if not inherited_typeface:
            continue
        variant_mapping = (font_variant_mapping or {}).get(inherited_typeface)
        if not variant_mapping:
            original_typeface = next(
                (
                    source
                    for source, replacement in font_mapping.items()
                    if replacement == inherited_typeface
                ),
                None,
            )
            if original_typeface:
                variant_mapping = (font_variant_mapping or {}).get(original_typeface)
                inherited_typeface = original_typeface
        if not variant_mapping:
            continue
        for child in paragraph:
            if child.tag not in run_tags:
                continue
            r_pr = child.find("a:rPr", PPT_NS)
            if r_pr is None:
                r_pr = ET.Element(f"{{{PPT_NS['a']}}}rPr")
                child.insert(0, r_pr)
            if _first_typeface(r_pr):
                continue
            variant = _font_style_variant(inherited_typeface, r_pr, [paragraph_default])
            replacement = variant_mapping.get(variant)
            if not replacement:
                continue
            latin = ET.SubElement(r_pr, f"{{{PPT_NS['a']}}}latin")
            latin.set("typeface", replacement)
            changed = True
    return changed


def _replace_fonts_in_pptx_xml(
    pptx_path: str,
    font_mapping: Dict[str, str],
    output_path: str,
    font_variant_mapping: Optional[Dict[str, Dict[str, str]]] = None,
) -> None:
    xml_prefixes = (
        "ppt/slides/",
        "ppt/slideLayouts/",
        "ppt/slideMasters/",
        "ppt/charts/",
    )
    with zipfile.ZipFile(pptx_path, "r") as src, zipfile.ZipFile(
        output_path, "w", compression=zipfile.ZIP_DEFLATED
    ) as dst:
        for info in src.infolist():
            data = src.read(info.filename)
            if info.filename.endswith(".xml") and info.filename.startswith(xml_prefixes):
                try:
                    root = ET.fromstring(data)
                    if _replace_fonts_in_xml_root(
                        root, font_mapping, font_variant_mapping
                    ):
                        data = ET.tostring(
                            root, encoding="utf-8", xml_declaration=True
                        )
                except Exception:
                    pass
            dst.writestr(info, data)


def replace_fonts_in_pptx(
    pptx_path: str,
    font_mapping: Dict[str, str],
    output_path: str,
    font_variant_mapping: Optional[Dict[str, Dict[str, str]]] = None,
) -> None:
    """
    Replace fonts in a PPTX file using python-pptx.

    Args:
        pptx_path: Path to input PPTX file
        font_mapping: Dictionary mapping old font names to new font names
        output_path: Path to save modified PPTX file
    """
    if font_variant_mapping:
        _replace_fonts_in_pptx_xml(
            pptx_path, font_mapping, output_path, font_variant_mapping
        )
        return

    if font_mapping:
        _replace_fonts_in_pptx_xml(pptx_path, font_mapping, output_path)
        return

    prs = Presentation(pptx_path)

    # Iterate through all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and run.font.name in font_mapping:
                            run.font.name = font_mapping[run.font.name]

            # Handle tables safely (python-pptx raises ValueError if non-table)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name and run.font.name in font_mapping:
                                    run.font.name = font_mapping[run.font.name]

    # Update slide layouts
    for slide_layout in prs.slide_layouts:
        for shape in slide_layout.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and run.font.name in font_mapping:
                            run.font.name = font_mapping[run.font.name]

    # Update slide masters
    for slide_master in prs.slide_masters:
        for shape in slide_master.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and run.font.name in font_mapping:
                            run.font.name = font_mapping[run.font.name]

    # Save the modified presentation
    prs.save(output_path)


def extract_font_from_eot(eot_path: Path) -> bytes:
    """Extract embedded font data from an EOT file."""
    with open(eot_path, "rb") as f:
        data = f.read()

    # EOT file structure:
    # - Header (variable length)
    # - Font family name (Unicode, null-terminated)
    # - Font style name (Unicode, null-terminated)
    # - Font version (Unicode, null-terminated)
    # - Font full name (Unicode, null-terminated)
    # - RootString (Unicode, null-terminated)
    # - Signature (4 bytes: "BSGP")
    # - Embedded font data (TTF/OTF) starts with "OTTO" or "ttcf"

    # Find the OpenType font signature - this marks the start of the embedded font
    # "OTTO" = OpenType with CFF (PostScript outlines)
    # "ttcf" = TrueType Collection
    # "\x00\x01\x00\x00" = TrueType with TrueType outlines
    otto_pos = data.find(b"OTTO")
    ttcf_pos = data.find(b"ttcf")
    ttf_pos = data.find(b"\x00\x01\x00\x00")

    font_start = -1
    if otto_pos != -1:
        font_start = otto_pos
    elif ttcf_pos != -1:
        font_start = ttcf_pos
    elif ttf_pos != -1:
        font_start = ttf_pos

    if font_start == -1:
        raise ValueError(
            "Could not find embedded font signature (OTTO/ttcf/TTF) in EOT file"
        )

    # Extract the embedded font from the found position to the end of file
    embedded_font = data[font_start:]

    return embedded_font


def get_font_details(path: str) -> FontDetail:
    """Extract detailed information from a font file."""
    font_path = Path(path)
    details = {
        "file": path,
        "size_bytes": font_path.stat().st_size,
        "error": None,
    }

    try:
        # Check if it's an EOT file
        is_eot = (
            font_path.suffix.lower() == ".fntdata" or font_path.suffix.lower() == ".eot"
        )

        if is_eot:
            # Extract embedded font from EOT
            try:
                embedded_font_data = extract_font_from_eot(font_path)
                # Create a temporary file to hold the extracted font
                with tempfile.NamedTemporaryFile(
                    delete=False, suffix=".ttf"
                ) as tmp_file:
                    tmp_file.write(embedded_font_data)
                    tmp_path = tmp_file.name

                try:
                    font = TTFont(tmp_path)
                finally:
                    # Clean up temp file
                    try:
                        os.unlink(tmp_path)
                    except Exception:
                        pass
            except Exception as e:
                # If extraction fails, try reading EOT metadata directly
                details["eot_extraction_error"] = str(e)
                # Fall through to try direct reading
                font = TTFont(str(font_path))
        else:
            # Try to open the font file directly
            font = TTFont(str(font_path))

        # Get font names from the 'name' table
        name_table = font.get("name")
        if name_table:
            names = {}
            for record in name_table.names:
                name_id = record.nameID
                platform_id = record.platformID
                # Prefer Unicode names (platformID 3) or Mac (platformID 1)
                if platform_id in (1, 3) or name_id not in names:
                    try:
                        name_str = (
                            record.toUnicode()
                            if hasattr(record, "toUnicode")
                            else str(record)
                        )
                        if name_str:
                            cleaned_name = _clean_font_metadata_string(name_str)
                            if cleaned_name:
                                names[name_id] = cleaned_name
                    except Exception:
                        pass

            # Map name IDs to readable names
            name_mapping = {
                1: "family_name",
                2: "subfamily_name",
                3: "unique_id",
                4: "full_name",
                5: "version",
                6: "postscript_name",
                7: "trademark",
                8: "manufacturer",
                9: "designer",
                10: "description",
                11: "vendor_url",
                12: "designer_url",
                13: "license",
                14: "license_url",
            }

            for name_id, key in name_mapping.items():
                if name_id in names:
                    details[key] = names[name_id]

        # Get OS/2 table for additional metrics
        os2_table = font.get("OS/2")
        if os2_table:
            details["weight_class"] = os2_table.usWeightClass
            details["width_class"] = os2_table.usWidthClass
            details["cap_height"] = getattr(os2_table, "sCapHeight", None)
            details["x_height"] = getattr(os2_table, "sxHeight", None)
            details["ascent"] = getattr(os2_table, "usWinAscent", None)
            details["descent"] = getattr(os2_table, "usWinDescent", None)

        # Get head table for font metrics
        head_table = font.get("head")
        if head_table:
            details["units_per_em"] = head_table.unitsPerEm
            details["created"] = head_table.created
            details["modified"] = head_table.modified

        # Get hhea table for horizontal metrics
        hhea_table = font.get("hhea")
        if hhea_table:
            details["ascender"] = hhea_table.ascent
            details["descender"] = hhea_table.descent
            details["line_gap"] = hhea_table.lineGap

        # Get number of glyphs
        if "cmap" in font:
            details["num_glyphs"] = len(font.getGlyphSet())

        # Get font format
        if hasattr(font, "sfntVersion"):
            details["format"] = _normalize_font_format(font.sfntVersion)

        font.close()

    except Exception as e:
        details["error"] = str(e)

    return FontDetail(**details)


def convert_eot_to_ttf(inp_path: str, out_dir: str) -> str:
    """
    Convert an EOT file to TTF format.

    Args:
        inp_path: Path to the input EOT file
        out_dir: Output directory where the converted font file will be saved

    Returns:
        Path to the converted font file
    """
    eot_path = Path(inp_path)
    out_dir_path = Path(out_dir)

    if not eot_path.exists():
        raise FileNotFoundError(f"EOT file not found: {eot_path}")

    # Create output directory if it doesn't exist
    out_dir_path.mkdir(parents=True, exist_ok=True)

    # Extract embedded font from EOT
    embedded_font_data = extract_font_from_eot(eot_path)

    # Determine the font format from the signature
    if embedded_font_data.startswith(b"OTTO"):
        font_format = "otf"
        default_ext = ".otf"
    elif embedded_font_data.startswith(b"ttcf"):
        font_format = "ttc"  # TrueType Collection
        default_ext = ".ttc"
    elif embedded_font_data.startswith(b"\x00\x01\x00\x00"):
        font_format = "ttf"
        default_ext = ".ttf"
    else:
        # Default to TTF if we can't determine
        font_format = "ttf"
        default_ext = ".ttf"

    # Construct output path in the output directory
    output_path = out_dir_path / f"{eot_path.stem}{default_ext}"

    # If the embedded font is OTF but output is requested as TTF,
    # attempt conversion using fonttools
    if font_format == "otf" and output_path.suffix.lower() == ".ttf":
        try:
            # Write to temp file first
            with tempfile.NamedTemporaryFile(delete=False, suffix=".otf") as tmp_file:
                tmp_file.write(embedded_font_data)
                tmp_otf_path = tmp_file.name

            try:
                # Open the OTF font
                font = TTFont(tmp_otf_path)

                # Try to convert CFF to TrueType outlines
                # This is a complex process - fonttools can't directly convert CFF to TTF
                # but we can try to save it and see if it works
                # Note: This may not work perfectly for all fonts
                font.flavor = None

                # Save as TTF (fonttools will attempt conversion)
                font.save(output_path)
                font.close()

                print(
                    "Note: Converted OTF (PostScript outlines) to TTF format. "
                    "Some glyph outlines may need manual adjustment."
                )
            finally:
                # Clean up temp file
                try:
                    os.unlink(tmp_otf_path)
                except Exception:
                    pass
        except Exception as e:
            # If conversion fails, save as OTF instead
            actual_output = output_path.with_suffix(".otf")
            with open(actual_output, "wb") as f:
                f.write(embedded_font_data)
            print(
                f"Warning: Could not convert OTF to TTF ({e}). "
                f"Saved as {actual_output.name} instead. "
                f"OTF to TTF conversion requires glyph outline recompilation."
            )
            return str(actual_output)
    else:
        # Write the extracted font directly to the output file
        with open(output_path, "wb") as f:
            f.write(embedded_font_data)

    return str(output_path)


_WEIGHT_KEYWORDS = {
    "thin": (
        "thin",
        "hairline",
    ),
    "extra_light": (
        "extra light",
        "extra-light",
        "extralight",
        "ultra light",
        "ultra-light",
        "ultralight",
    ),
    "light": ("light",),
    "regular": (
        "regular",
        "normal",
        "book",
    ),
    "medium": ("medium",),
    "semibold": (
        "semi bold",
        "semi-bold",
        "semibold",
        "demi bold",
        "demi-bold",
        "demibold",
    ),
    "bold": ("bold",),
    "extra_bold": (
        "extra bold",
        "extra-bold",
        "extrabold",
        "ultra bold",
        "ultra-bold",
        "ultrabold",
    ),
    "black": (
        "black",
        "heavy",
    ),
    "extra_black": (
        "extra black",
        "extra-black",
        "extrablack",
        "ultra black",
        "ultra-black",
        "ultrablack",
        "super black",
        "super-black",
        "superblack",
    ),
}

_STYLE_KEYWORDS = ("italic", "oblique")

_WEIGHT_CLASS_BUCKETS = (
    ("thin", 0, 149),
    ("extra_light", 150, 249),
    ("light", 250, 349),
    ("regular", 350, 449),
    ("medium", 450, 549),
    ("semibold", 550, 649),
    ("bold", 650, 749),
    ("extra_bold", 750, 849),
    ("black", 850, 925),
    ("extra_black", 926, 1000),
)


def _normalize_text(value: Optional[str]) -> str:
    if not value:
        return ""
    lowered = value.lower()
    lowered = re.sub(r"[^a-z0-9]+", " ", lowered)
    return re.sub(r"\s+", " ", lowered).strip()


def _normalize_compact(value: Optional[str]) -> str:
    if not value:
        return ""
    return re.sub(r"[^a-z0-9]+", "", value.lower())


@lru_cache(maxsize=1)
def _get_weight_keyword_index():
    entries = []
    for canonical, phrases in _WEIGHT_KEYWORDS.items():
        for phrase in phrases:
            normalized = _normalize_text(phrase)
            compact = _normalize_compact(phrase)
            if normalized:
                entries.append((normalized, compact, canonical))
    entries.sort(key=lambda item: len(item[0]), reverse=True)
    return tuple(entries)


@lru_cache(maxsize=1)
def _get_removal_keywords():
    keywords = set()
    for normalized, _, _ in _get_weight_keyword_index():
        if normalized:
            keywords.add(normalized)
    for style in _STYLE_KEYWORDS:
        normalized_style = _normalize_text(style)
        if normalized_style:
            keywords.add(normalized_style)
    return tuple(sorted(keywords, key=len, reverse=True))


def _family_key(value: Optional[str]) -> str:
    normalized = _normalize_text(value)
    if not normalized:
        return ""
    cleaned = normalized
    for keyword in _get_removal_keywords():
        pattern = r"\b" + re.escape(keyword) + r"\b"
        cleaned = re.sub(pattern, " ", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    target = cleaned if cleaned else normalized
    return re.sub(r"\s+", "", target)


def _extract_weight_from_name(value: Optional[str]) -> Optional[str]:
    normalized = _normalize_text(value)
    if not normalized:
        return None
    compact = _normalize_compact(value)
    padded = f" {normalized} "
    for phrase_norm, phrase_compact, canonical in _get_weight_keyword_index():
        if not phrase_norm:
            continue
        if f" {phrase_norm} " in padded:
            return canonical
        if phrase_compact and phrase_compact in compact:
            return canonical
    return None


def _weight_from_class(weight_class: Optional[int]) -> Optional[str]:
    if weight_class is None:
        return None
    for canonical, lower, upper in _WEIGHT_CLASS_BUCKETS:
        if lower <= weight_class <= upper:
            return canonical
    return None


def _extract_weight_from_detail(font_detail: FontDetail) -> Optional[str]:
    weight = _weight_from_class(font_detail.weight_class)
    if weight:
        return weight
    for candidate in (
        font_detail.subfamily_name,
        font_detail.full_name,
        font_detail.postscript_name,
    ):
        weight = _extract_weight_from_name(candidate)
        if weight:
            return weight
    return None


def _weight_value_from_canonical(weight_key: Optional[str]) -> Optional[int]:
    if not weight_key:
        return None
    for canonical, lower, upper in _WEIGHT_CLASS_BUCKETS:
        if canonical == weight_key:
            return (lower + upper) // 2
    return None


def _weight_value_from_detail(font_detail: FontDetail) -> Optional[int]:
    canonical = _extract_weight_from_detail(font_detail)
    midpoint = _weight_value_from_canonical(canonical)
    if midpoint is not None:
        return midpoint
    if font_detail.weight_class is not None:
        for _, lower, upper in _WEIGHT_CLASS_BUCKETS:
            if lower <= font_detail.weight_class <= upper:
                return (lower + upper) // 2
    return None


def get_index_of_matching_font_detail_or_none(
    font_name: str, font_details: Sequence[FontDetail]
) -> Optional[int]:
    """
    Return the index of the font detail that best matches the provided font name.
    Family equality must match. If the requested weight is unspecified we treat it
    as a regular weight, but we still fall back to the closest match when no exact
    or regular-weight match can be found.
    """
    if not font_name or not font_details:
        return None

    family_key = _family_key(font_name)
    if not family_key:
        return None

    font_weight = _extract_weight_from_name(font_name)
    expected_weight = font_weight or "regular"
    expected_weight_value = _weight_value_from_canonical(expected_weight) or 400

    best_index: Optional[int] = None
    best_score = -1
    fallback_index: Optional[int] = None
    fallback_diff = float("inf")

    for index, font_detail in enumerate(font_details):
        if not font_detail:
            continue

        detail_keys = set()
        for value in (
            font_detail.full_name,
            font_detail.postscript_name,
            font_detail.family_name,
            font_detail.subfamily_name,
        ):
            key = _family_key(value)
            if key:
                detail_keys.add(key)

        if not detail_keys or family_key not in detail_keys:
            continue

        detail_weight = _extract_weight_from_detail(font_detail)
        detail_weight_value = (
            _weight_value_from_detail(font_detail) or expected_weight_value
        )

        score = 1
        if detail_weight == expected_weight:
            score = 3
        elif detail_weight is None and expected_weight == "regular":
            score = 2

        if score > best_score:
            best_index = index
            best_score = score

        diff = abs(detail_weight_value - expected_weight_value)
        if diff < fallback_diff:
            fallback_index = index
            fallback_diff = diff

    if best_score >= 2:
        return best_index
    return fallback_index


def extract_font_name_from_file(file_path: str) -> str:
    """Extract the canonical font family name from a font file."""
    filename = os.path.basename(file_path)
    try:
        font = TTFont(file_path)
        if "name" in font:
            name_table = font["name"]
            for name_id in [1, 4, 6]:
                for record in name_table.names:
                    if record.nameID == name_id:
                        if record.langID == 0x409 or record.langID == 0:
                            font_name = record.toUnicode().strip()
                            if font_name:
                                font.close()
                                return font_name
            for record in name_table.names:
                if record.nameID == 1:
                    font_name = record.toUnicode().strip()
                    if font_name:
                        font.close()
                        return font_name
        font.close()
    except Exception as exc:
        print(f"[FONT DEBUG] Error reading font metadata for {filename}: {exc}")
    base_name = os.path.splitext(filename)[0]
    if "_" in filename and len(filename.split("_")[-1].split(".")[0]) == 8:
        parts = filename.split("_")
        if len(parts) > 1:
            return "_".join(parts[:-1])
    return base_name
